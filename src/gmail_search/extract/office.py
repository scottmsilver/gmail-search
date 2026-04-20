"""Extract text from Office documents (docx, xlsx, csv)."""

import logging
from pathlib import Path
from typing import Any

from gmail_search.extract import ExtractResult

logger = logging.getLogger(__name__)


def extract_docx(file_path: Path, config: dict[str, Any]) -> ExtractResult | None:
    """Extract paragraph text from Word documents."""
    from docx import Document

    try:
        doc = Document(str(file_path))
    except Exception as e:
        logger.warning(f"Failed to open docx {file_path}: {e}")
        return None

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return None

    return ExtractResult(text="\n\n".join(paragraphs))


def extract_xlsx(file_path: Path, config: dict[str, Any]) -> ExtractResult | None:
    """Extract cell text from Excel spreadsheets — handles both the
    modern .xlsx (OOXML, via openpyxl) and legacy .xls (binary
    BIFF, via xlrd). Dispatches on file magic, not extension, so
    mime-lies are tolerated.
    """
    # .xls starts with the OLE compound-file marker D0 CF 11 E0;
    # .xlsx is a zip archive starting with PK. Cheap two-byte sniff.
    try:
        header = file_path.read_bytes()[:8]
    except Exception as e:
        logger.warning(f"Failed to read xls/xlsx header {file_path}: {e}")
        return None
    if header.startswith(b"\xd0\xcf\x11\xe0"):
        return _extract_xls_legacy(file_path, config)

    from openpyxl import load_workbook

    try:
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
    except Exception as e:
        logger.warning(f"Failed to open xlsx {file_path}: {e}")
        return None

    parts: list[str] = []
    max_rows = config.get("max_xlsx_rows", 500)

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows_text: list[str] = []
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            if i >= max_rows:
                break
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows_text.append(" | ".join(cells))

        if rows_text:
            parts.append(f"Sheet: {sheet_name}\n" + "\n".join(rows_text))

    wb.close()

    if not parts:
        return None

    return ExtractResult(text="\n\n".join(parts))


def _extract_xls_legacy(file_path: Path, config: dict[str, Any]) -> ExtractResult | None:
    """Legacy binary .xls (OLE / BIFF) via xlrd. openpyxl does NOT
    read this format — it'd raise InvalidFileException. xlrd 2.0.x
    dropped xlsx on purpose, so the two libs are complementary.
    """
    try:
        import xlrd
    except ImportError:
        logger.warning("xlrd not installed; skipping %s", file_path)
        return None

    try:
        book = xlrd.open_workbook(str(file_path))
    except Exception as e:
        logger.warning(f"Failed to open xls {file_path}: {e}")
        return None

    max_rows = config.get("max_xlsx_rows", 500)
    parts: list[str] = []
    for sheet in book.sheets():
        rows_text: list[str] = []
        for r in range(min(sheet.nrows, max_rows)):
            cells = [
                str(sheet.cell_value(r, c)) for c in range(sheet.ncols) if sheet.cell_value(r, c) not in ("", None)
            ]
            if cells:
                rows_text.append(" | ".join(cells))
        if rows_text:
            parts.append(f"Sheet: {sheet.name}\n" + "\n".join(rows_text))
    if not parts:
        return None
    return ExtractResult(text="\n\n".join(parts))


def extract_csv(file_path: Path, config: dict[str, Any]) -> ExtractResult | None:
    """Extract text from CSV files."""
    max_rows = config.get("max_xlsx_rows", 500)

    try:
        text = file_path.read_text(errors="replace")
        lines = text.strip().split("\n")[:max_rows]
        if not lines:
            return None
        return ExtractResult(text="\n".join(lines))
    except Exception as e:
        logger.warning(f"Failed to read CSV {file_path}: {e}")
        return None


def extract_pptx(file_path: Path, config: dict[str, Any]) -> ExtractResult | None:
    """Extract slide text from PowerPoint .pptx.

    Concatenates every text frame on every slide. Slide boundaries
    are preserved as "--- Slide N ---" markers so the summarizer can
    still see ordering.
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed; skipping %s", file_path)
        return None

    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        logger.warning(f"Failed to open pptx {file_path}: {e}")
        return None

    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in p.runs).strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))
    if not parts:
        return None
    return ExtractResult(text="\n\n".join(parts))


def extract_doc(file_path: Path, config: dict[str, Any]) -> ExtractResult | None:
    """Legacy Word .doc (binary BIFF / OLE) via LibreOffice's
    `soffice --headless --convert-to txt` subprocess. LibreOffice is
    the canonical tool for .doc → plain-text conversion, maintained
    to handle every Word quirk (tables, nested lists, embedded
    objects) we'd otherwise reimplement.

    Returns None if LibreOffice isn't installed; we expose the
    missing-tool via a warning so the operator can install it.
    """
    import shutil
    import subprocess
    import tempfile

    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if soffice is None:
        logger.warning("libreoffice not installed; cannot extract %s", file_path)
        return None

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, str(file_path)],
                capture_output=True,
                timeout=60,
                check=True,
            )
        except subprocess.CalledProcessError as e:
            logger.warning(f"libreoffice failed on {file_path}: {e.stderr.decode(errors='replace')[:200]}")
            return None
        except subprocess.TimeoutExpired:
            logger.warning(f"libreoffice timed out on {file_path}")
            return None

        out_path = Path(tmpdir) / (file_path.stem + ".txt")
        if not out_path.exists():
            return None
        text = out_path.read_text(errors="replace").strip()

    if len(text) < 80:
        return None
    return ExtractResult(text=text[:50000])
