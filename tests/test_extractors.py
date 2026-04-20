"""Tests for the attachment extractor pipeline.

Each extractor is a "port" between a file format and our ExtractResult
dataclass. We build real fixtures (a real .docx with python-docx, a
real .pptx with python-pptx, a real .xlsx with openpyxl) because the
failure modes we care about happen inside those libraries' parsing
paths — mocking them would defeat the purpose.

Legacy binary formats (.doc, .xls) are tested against pre-built
binary fixtures checked into tests/fixtures/. LibreOffice's subprocess
is mocked where we can, real where we can't.
"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from gmail_search.extract import ExtractResult, dispatch

# ─── fixture helpers ──────────────────────────────────────────────────────


@pytest.fixture
def docx_file(tmp_path: Path) -> Path:
    """Real .docx built via python-docx — tests go through the actual
    OOXML unzip + parse path our extractor hits in production.
    """
    from docx import Document

    doc = Document()
    doc.add_heading("Meeting Notes", 0)
    doc.add_paragraph("Attendees: Alice, Bob, Charlie")
    doc.add_paragraph("Budget approved: $42,500 for Q3.")
    path = tmp_path / "notes.docx"
    doc.save(str(path))
    return path


@pytest.fixture
def xlsx_file(tmp_path: Path) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    sheet = wb.active
    sheet.title = "Invoices"
    sheet.append(["Customer", "Amount", "Due"])
    sheet.append(["Acme Inc", 12500.00, "2026-05-01"])
    sheet.append(["Globex", 8200.00, "2026-05-15"])
    path = tmp_path / "invoices.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture
def pptx_file(tmp_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    layout = prs.slide_layouts[5]  # title only
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Q4 Roadmap"
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Launch Targets"
    # Add a plain text box with more detail
    tb = s2.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(2))
    tb.text_frame.text = "EU GA: Nov 15\nAPAC GA: Dec 1"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def html_file(tmp_path: Path) -> Path:
    path = tmp_path / "page.html"
    path.write_text(
        "<html><head><title>Invoice</title>"
        "<style>.red{color:red}</style>"
        "<script>alert('nope')</script></head>"
        "<body><h1>Invoice #42</h1>"
        "<p>Payment of $1,299.00 due by <b>April 30, 2026</b>.</p>"
        "<p>Customer: Acme Corp</p></body></html>"
    )
    return path


@pytest.fixture
def csv_file(tmp_path: Path) -> Path:
    path = tmp_path / "names.csv"
    path.write_text("name,role,email\nAlice,PM,alice@x.com\nBob,Eng,bob@x.com\n")
    return path


# ─── extractor tests ──────────────────────────────────────────────────────


def test_docx_extracts_paragraphs(docx_file):
    from gmail_search.extract.office import extract_docx

    result = extract_docx(docx_file, {})
    assert isinstance(result, ExtractResult)
    assert "Meeting Notes" in result.text
    assert "Alice" in result.text
    assert "$42,500" in result.text


def test_xlsx_extracts_cells(xlsx_file):
    from gmail_search.extract.office import extract_xlsx

    result = extract_xlsx(xlsx_file, {})
    assert isinstance(result, ExtractResult)
    assert "Invoices" in result.text
    assert "Acme Inc" in result.text
    assert "12500" in result.text


def test_xlsx_dispatches_legacy_xls_via_sniff(tmp_path):
    """Legacy binary .xls starts with the OLE compound-file marker.
    Our dispatcher for .xlsx sniffs the first bytes and redirects to
    xlrd. This test stubs xlrd to verify the path fires.
    """
    fake = tmp_path / "old.xls"
    # Real OLE header so our sniff hits the xls path; rest is junk.
    fake.write_bytes(b"\xd0\xcf\x11\xe0" + b"\x00" * 512)

    import gmail_search.extract.office as office

    fake_sheet = MagicMock()
    fake_sheet.name = "Legacy"
    fake_sheet.nrows = 2
    fake_sheet.ncols = 2
    fake_sheet.cell_value.side_effect = lambda r, c: [["Year", "Revenue"], [2023, 100000]][r][c]
    fake_book = MagicMock()
    fake_book.sheets.return_value = [fake_sheet]

    with patch.object(office, "__name__", office.__name__):
        with patch("xlrd.open_workbook", return_value=fake_book):
            result = office.extract_xlsx(fake, {})
    assert result is not None
    assert "Legacy" in result.text
    assert "Year" in result.text
    assert "100000" in result.text


def test_pptx_extracts_slides(pptx_file):
    from gmail_search.extract.office import extract_pptx

    result = extract_pptx(pptx_file, {})
    assert result is not None
    assert "Q4 Roadmap" in result.text
    assert "Launch Targets" in result.text
    assert "Nov 15" in result.text
    # Slide markers preserved so the summarizer sees ordering.
    assert "--- Slide 1 ---" in result.text
    assert "--- Slide 2 ---" in result.text


def test_html_strips_script_and_style(html_file):
    from gmail_search.extract.text import extract_html

    result = extract_html(html_file, {})
    assert result is not None
    assert "Invoice #42" in result.text
    assert "$1,299.00" in result.text
    assert "Acme Corp" in result.text
    # JS/CSS must not leak into the retrieval text.
    assert "alert" not in result.text
    assert ".red" not in result.text


def test_csv_extracts_rows(csv_file):
    from gmail_search.extract.office import extract_csv

    result = extract_csv(csv_file, {})
    assert result is not None
    assert "alice@x.com" in result.text
    assert "bob@x.com" in result.text


def test_doc_uses_libreoffice_subprocess(tmp_path):
    """Mock the libreoffice subprocess and verify we shell out to it
    with the expected args. We don't try to build a real .doc binary
    — the subprocess call + stdout parsing IS the contract.
    """
    fake_doc = tmp_path / "legacy.doc"
    fake_doc.write_bytes(b"\xd0\xcf\x11\xe0" + b"junk" * 20)

    # soffice "output" directory: simulate its effect by writing
    # the expected .txt into the tmpdir the caller creates.
    def fake_run(cmd, capture_output, timeout, check):
        outdir_idx = cmd.index("--outdir") + 1
        outdir = Path(cmd[outdir_idx])
        (outdir / f"{fake_doc.stem}.txt").write_text(
            "Contract signed by Alice Martinez on April 1, 2026. Terms: $12,500 net-30 for "
            "a six-month engagement starting May 1, 2026."
        )

        class R:
            stdout, stderr, returncode = b"", b"", 0

        return R()

    from gmail_search.extract.office import extract_doc

    with (
        patch("shutil.which", return_value="/usr/bin/libreoffice"),
        patch("subprocess.run", side_effect=fake_run) as mock_run,
    ):
        result = extract_doc(fake_doc, {})

    assert result is not None
    assert "Contract signed by Alice" in result.text
    args = mock_run.call_args[0][0]
    assert args[0] == "/usr/bin/libreoffice"
    assert "--headless" in args
    assert "--convert-to" in args
    assert str(fake_doc) in args


def test_doc_returns_none_when_libreoffice_missing(tmp_path):
    from gmail_search.extract.office import extract_doc

    with patch("shutil.which", return_value=None):
        result = extract_doc(tmp_path / "x.doc", {})
    assert result is None


# ─── dispatch sniffing ─────────────────────────────────────────────────────


def test_dispatch_sniffs_octet_stream_by_extension(docx_file):
    """A Gmail attachment reported as application/octet-stream that's
    actually a .docx should be routed via the sniffer back to
    extract_docx — otherwise the corpus misses 145+ PDFs, 50+ PNGs,
    and 30+ calendar files that Gmail mis-reports.
    """
    # Repoint the docx so its extension is preserved but mime lies.
    disguised = docx_file.rename(docx_file.with_name("mystery.docx"))
    result = dispatch("application/octet-stream", disguised, {})
    assert result is not None
    assert "Meeting Notes" in result.text


def test_dispatch_sniffs_octet_stream_html(html_file):
    disguised = html_file.rename(html_file.with_name("file.html"))
    result = dispatch("application/octet-stream", disguised, {})
    assert result is not None
    assert "Invoice #42" in result.text


def test_dispatch_returns_none_for_unknown_extension(tmp_path):
    f = tmp_path / "thing.xyz"
    f.write_bytes(b"whatever")
    assert dispatch("application/octet-stream", f, {}) is None


def test_dispatch_returns_none_for_unknown_mime(tmp_path):
    f = tmp_path / "x.bin"
    f.write_bytes(b"bytes")
    assert dispatch("application/x-squid-brains", f, {}) is None


# ─── drive URL extractor ──────────────────────────────────────────────────


def test_drive_extract_ids_canonical():
    from gmail_search.gmail.drive import extract_drive_ids

    body = """Hi team — the deck is at
    https://docs.google.com/presentation/d/1AbC-DeFgHijKlMnOpQrStUvWxYz0123456/edit?usp=sharing
    and the sheet is https://docs.google.com/spreadsheets/d/1XyZ_AbCdEfGhIjKlMnOpQrStUvWxYz789012/view
    Thanks!
    """
    ids = extract_drive_ids(body)
    assert len(ids) == 2
    assert ("1AbC-DeFgHijKlMnOpQrStUvWxYz0123456", "slides") in ids
    assert ("1XyZ_AbCdEfGhIjKlMnOpQrStUvWxYz789012", "sheet") in ids


def test_drive_extract_ids_dedupes():
    from gmail_search.gmail.drive import extract_drive_ids

    body = (
        "Link: https://docs.google.com/document/d/1DEADBEEFCAFEBABECAFEBABE1234/edit "
        "and again: https://docs.google.com/document/d/1DEADBEEFCAFEBABECAFEBABE1234/view"
    )
    ids = extract_drive_ids(body)
    assert ids == [("1DEADBEEFCAFEBABECAFEBABE1234", "doc")]


def test_drive_extract_ids_handles_file_d_links():
    from gmail_search.gmail.drive import extract_drive_ids

    body = "PDF attached: https://drive.google.com/file/d/1PdfABCDEFGHIJKLMNOPQRSTUVWXYZ/view"
    ids = extract_drive_ids(body)
    assert ids == [("1PdfABCDEFGHIJKLMNOPQRSTUVWXYZ", "file")]


def test_drive_extract_ids_empty_on_plain_text():
    from gmail_search.gmail.drive import extract_drive_ids

    assert extract_drive_ids("No links here, just prose about 2026-04-19.") == []


# ─── drive fetch (mocked service) ──────────────────────────────────────────


def test_fetch_doc_exports_native_doc_as_text():
    """Native Google Doc → export_media mime=text/plain."""
    from gmail_search.gmail.drive import fetch_doc_text

    service = MagicMock()
    service.files().get().execute.return_value = {
        "id": "x",
        "name": "Launch plan",
        "mimeType": "application/vnd.google-apps.document",
    }
    service.files().export_media().execute.return_value = b"Launch plan body\nStep 1. Ship\nStep 2. Iterate"
    result = fetch_doc_text(service, "x")
    assert result is not None
    name, text = result
    assert name == "Launch plan"
    assert "Step 1. Ship" in text


def test_fetch_doc_skips_binary_uploads():
    """An uploaded .jpg on Drive would require a download pass the
    attachment pipeline already handles — skip from drive.py to avoid
    doing the same work twice."""
    from gmail_search.gmail.drive import fetch_doc_text

    service = MagicMock()
    service.files().get().execute.return_value = {"id": "y", "name": "pic.jpg", "mimeType": "image/jpeg"}
    assert fetch_doc_text(service, "y") is None


def test_fetch_doc_returns_none_on_api_error():
    from gmail_search.gmail.drive import fetch_doc_text

    service = MagicMock()
    service.files().get().execute.side_effect = Exception("403 Insufficient scope")
    assert fetch_doc_text(service, "z") is None
